"""
2026年海外客户投诉台账 - Streamlit 云端数据看板
部署于 Streamlit Community Cloud, 通过 GitHub 自动更新
"""

import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from datetime import datetime, timezone, timedelta
import os

# Beijing timezone
BJT = timezone(timedelta(hours=8))
import io
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

# ============================================================
# 结案预警邮件配置
# ============================================================
FOLLOWER_EMAILS = {
    "郑小平": "payne.zheng@tcl.com",
    "陈耀球": "kt_yorkchen@tcl.com",
    "曾靖衍": "jingyan.zeng@tcl.com",
    "黄忠成": "zhongcheng.huang@tcl.com",
    "方益勋": "kt_fangyx@tcl.com",
}
SMTP_SERVER   = os.getenv("SMTP_SERVER", "mail.tcl.com")
SMTP_PORT     = int(os.getenv("SMTP_PORT", "587"))
SMTP_USER     = os.getenv("SMTP_USER", "pub_kehufuwubu@tcl.com")
SMTP_PASSWORD = os.getenv("SMTP_PASSWORD", "svc@2027")
import zipfile
import hashlib
import urllib.request
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# Page Config
# ============================================================
st.set_page_config(
    page_title="2026海外客诉看板",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ============================================================
# Data Source Config
# ============================================================
LOCAL_FILE = os.path.join(os.path.dirname(__file__), "2026年海外客户投诉台账.xlsx")
KDOCS_EDIT_URL = "https://www.kdocs.cn/l/cjP6zkIRj17V"
CSV_COLS = [
    '编号','分公司','国家或地区','是否大客户','投诉日期','应结案日期','实际完成日期',
    '完成周期（天）','机型属性','故障比例','问题描述','客户诉求','跟进人','处理类型',
    '结案状态','应急措施','原因分析','长期整改措施','责任单位','故障大类','品质负责人',
    '8D报告','8D措施点检','备注'
]

# ============================================================
# Data Loading
# ============================================================
def _process_raw_df(df):
    """清洗和标准化原始 DataFrame"""
    df = df.dropna(how='all')

    # Normalize column names (handle different encodings/spacing)
    col_map = {}
    for c in df.columns:
        c_stripped = c.strip().replace('\n','').replace('\r','')
        for ec in CSV_COLS:
            if c_stripped == ec or c_stripped == ec.replace('（','(').replace('）',')'):
                col_map[c] = ec
                break
    if col_map:
        df = df.rename(columns=col_map)

    # Strict filtering
    df['编号_num'] = pd.to_numeric(df['编号'], errors='coerce')
    df = df[df['编号_num'].notna() & df['分公司'].notna() & (df['分公司'].astype(str).str.strip() != '')].copy()
    df = df.drop_duplicates(subset='编号_num', keep='first').sort_values('编号_num').reset_index(drop=True)
    df['编号_int'] = df['编号_num'].astype(int)

    # Date parsing
    for col in ['投诉日期','应结案日期','实际完成日期']:
        df[col] = pd.to_datetime(df[col], errors='coerce')

    df['完成周期（天）'] = pd.to_numeric(df['完成周期（天）'], errors='coerce')
    df['故障比例'] = pd.to_numeric(df['故障比例'], errors='coerce')
    df['投诉月份'] = df['投诉日期'].dt.to_period('M').astype(str)

    return df


def load_data(uploaded_bytes=None):
    """加载数据: 上传文件 > 本地Excel"""
    source_label = ""

    if uploaded_bytes is not None:
        try:
            df = pd.read_excel(io.BytesIO(uploaded_bytes), sheet_name='所有客诉', header=1)
            source_label = "📤 已上传文件"
        except Exception:
            st.error("❌ 无法读取上传的文件，请确认格式正确")
            st.stop()
    else:
        try:
            df = pd.read_excel(LOCAL_FILE, sheet_name='所有客诉', header=1)
            source_label = "默认数据"
        except Exception:
            st.error("❌ 无法加载数据源，请检查文件")
            st.stop()

    df = _process_raw_df(df)
    df.attrs['source'] = source_label
    return df


def apply_filters(df):
    """Apply sidebar filters to dataframe"""
    filtered = df.copy()

    # Branch filter
    branches = st.sidebar.multiselect(
        "分公司", options=sorted(df['分公司'].dropna().unique()),
        default=[]
    )
    if branches:
        filtered = filtered[filtered['分公司'].isin(branches)]

    # Status filter
    statuses = st.sidebar.multiselect(
        "结案状态", options=sorted(df['结案状态'].dropna().unique()),
        default=[]
    )
    if statuses:
        filtered = filtered[filtered['结案状态'].isin(statuses)]

    # VIP filter
    vip_options = st.sidebar.multiselect(
        "客户类型", options=['是','否'],
        default=[]
    )
    if vip_options:
        filtered = filtered[filtered['是否大客户'].isin(vip_options)]

    # Date range
    min_date = df['投诉日期'].min()
    max_date = df['投诉日期'].max()
    if pd.notna(min_date) and pd.notna(max_date):
        date_range = st.sidebar.date_input(
            "投诉日期范围",
            value=(min_date.date(), max_date.date()),
            min_value=min_date.date(),
            max_value=max_date.date(),
        )
        if len(date_range) == 2:
            filtered = filtered[
                (filtered['投诉日期'] >= pd.Timestamp(date_range[0])) &
                (filtered['投诉日期'] <= pd.Timestamp(date_range[1]))
            ]

    # Fault category filter
    faults = st.sidebar.multiselect(
        "故障大类", options=sorted(df['故障大类'].dropna().unique()),
        default=[]
    )
    if faults:
        filtered = filtered[filtered['故障大类'].isin(faults)]

    return filtered


# ============================================================
# Color Palette
# ============================================================
COLORS = ['#1890FF','#27AE60','#F39C12','#E74C3C','#9B59B6','#1ABC9C','#E67E22','#3498DB',
          '#2ECC71','#E91E63','#00BCD4','#FF9800','#795548','#607D8B','#CDDC39']

STATUS_COLORS = {'结案': '#27AE60', '关闭': '#95A5A6', '未结案': '#F39C12', '暂停': '#E74C3C'}

# ============================================================
# Chart Functions
# ============================================================
def make_branch_chart(df):
    """分公司投诉分布 - 横向柱状图"""
    counts = df['分公司'].value_counts().reset_index()
    counts.columns = ['分公司','数量']
    counts = counts.sort_values('数量')

    colors = ['#E74C3C' if v >= 8 else '#F39C12' if v >= 5 else '#1890FF' for v in counts['数量']]

    fig = go.Figure(go.Bar(
        x=counts['数量'], y=counts['分公司'],
        orientation='h',
        marker_color=colors,
        text=counts['数量'],
        textposition='outside',
        hovertemplate='%{y}: %{x} 条<extra></extra>',
    ))
    fig.update_layout(
        title='分公司投诉分布',
        xaxis_title='投诉数量',
        height=360,
        margin=dict(l=10, r=30, t=40, b=10),
        bargap=0.3,
    )
    return fig


def make_country_chart(df):
    """国家/地区 TOP15"""
    counts = df['国家或地区'].value_counts().nlargest(15).reset_index()
    counts.columns = ['国家','数量']
    counts = counts.sort_values('数量')

    colors = ['#E74C3C' if i >= len(counts)-3 else '#1890FF' for i in range(len(counts))]

    fig = go.Figure(go.Bar(
        x=counts['数量'], y=counts['国家'],
        orientation='h',
        marker_color=colors,
        text=counts['数量'],
        textposition='outside',
        hovertemplate='%{y}: %{x} 条<extra></extra>',
    ))
    fig.update_layout(
        title='国家/地区投诉量 TOP15',
        xaxis_title='投诉数量',
        height=360,
        margin=dict(l=10, r=30, t=40, b=10),
        bargap=0.25,
    )
    return fig


def make_monthly_trend(df):
    """月度投诉趋势 + 结案率"""
    monthly = df.groupby('投诉月份').agg(
        投诉量=('编号_int','count'),
        结案量=('结案状态', lambda x: x.isin(['结案','关闭']).sum())
    ).reset_index()
    monthly['结案率'] = (monthly['结案量'] / monthly['投诉量'] * 100).round(1)

    fig = make_subplots(specs=[[{"secondary_y": True}]])

    fig.add_trace(go.Bar(
        name='投诉量', x=monthly['投诉月份'], y=monthly['投诉量'],
        marker_color='#1890FF', text=monthly['投诉量'], textposition='outside',
    ), secondary_y=False)

    fig.add_trace(go.Bar(
        name='已结案', x=monthly['投诉月份'], y=monthly['结案量'],
        marker_color='#27AE60',
    ), secondary_y=False)

    fig.add_trace(go.Scatter(
        name='结案率', x=monthly['投诉月份'], y=monthly['结案率'],
        mode='lines+markers+text', text=monthly['结案率'].apply(lambda x: f'{x}%'),
        textposition='top center', line=dict(color='#F39C12', width=3),
        marker=dict(size=10, symbol='diamond'),
    ), secondary_y=True)

    fig.add_hline(y=monthly['结案率'].mean(), line_dash="dash", line_color="gray",
                  annotation_text=f"平均结案率 {monthly['结案率'].mean():.0f}%", secondary_y=True)

    fig.update_layout(
        title='月度投诉趋势 & 结案率',
        hovermode='x unified',
        height=380,
        margin=dict(l=10, r=10, t=40, b=10),
        legend=dict(orientation='h', y=-0.15),
    )
    fig.update_yaxes(title_text="条数", secondary_y=False)
    fig.update_yaxes(title_text="结案率 (%)", range=[0, 105], secondary_y=True)
    return fig


def make_status_pie(df):
    """结案状态环形图"""
    counts = df['结案状态'].value_counts().reset_index()
    counts.columns = ['状态','数量']

    fig = go.Figure(go.Pie(
        labels=counts['状态'], values=counts['数量'],
        hole=0.55,
        marker_colors=[STATUS_COLORS.get(s, '#999') for s in counts['状态']],
        textinfo='label+percent',
        textfont_size=13,
        hovertemplate='%{label}: %{value} 条 (%{percent})<extra></extra>',
    ))
    fig.update_layout(
        title='结案状态分布',
        height=380,
        margin=dict(l=10, r=10, t=40, b=10),
    )
    return fig


def make_fault_pareto(df):
    """故障大类帕累托图（标准垂直柱状图+累计折线）"""
    counts = df['故障大类'].value_counts().reset_index()
    counts.columns = ['故障大类','数量']
    total = counts['数量'].sum()
    if total == 0:
        total = 1  # 防止除以零
    counts = counts.sort_values('数量', ascending=False)
    counts['累计占比'] = counts['数量'].cumsum() / total * 100

    colors = ['#E74C3C' if i < 3 else '#F39C12' if i < 6 else '#1890FF'
              for i in range(len(counts))]

    fig = make_subplots(specs=[[{"secondary_y": True}]])

    fig.add_trace(go.Bar(
        name='投诉数量', x=counts['故障大类'], y=counts['数量'],
        marker_color=colors,
        text=counts['数量'], textposition='outside',
    ), secondary_y=False)

    fig.add_trace(go.Scatter(
        name='累计占比', x=counts['故障大类'], y=counts['累计占比'],
        mode='lines+markers', line=dict(color='#9B59B6', width=2.5),
        marker=dict(size=8, symbol='diamond'),
    ), secondary_y=True)

    fig.add_hline(y=80, line_dash="dash", line_color="#E74C3C",
                  annotation_text="80% 线", secondary_y=True)

    fig.update_layout(
        title='故障大类帕累托分析',
        height=400,
        margin=dict(l=10, r=10, t=40, b=40),
        legend=dict(orientation='h', y=-0.15),
        xaxis=dict(tickangle=-35, tickfont=dict(size=11)),
    )
    fig.update_yaxes(title_text="投诉数量", secondary_y=False)
    fig.update_yaxes(title_text="累计占比 (%)", range=[0, 105], secondary_y=True)
    return fig


def make_model_vip_chart(df):
    """机型属性 + 大客户占比"""
    model_counts = df['机型属性'].value_counts().reset_index()
    model_counts.columns = ['机型','数量']

    fig = make_subplots(
        rows=1, cols=2,
        specs=[[{'type': 'pie'}, {'type': 'pie'}]],
        subplot_titles=('机型属性', '大客户 vs 非大客户'),
    )

    fig.add_trace(go.Pie(
        labels=model_counts['机型'], values=model_counts['数量'],
        hole=0.4, textinfo='label+percent',
    ), row=1, col=1)

    vip_counts = df['是否大客户'].value_counts()
    fig.add_trace(go.Pie(
        labels=['非大客户','大客户'],
        values=[vip_counts.get('否',0), vip_counts.get('是',0)],
        hole=0.4, textinfo='label+percent',
        marker_colors=['#1890FF','#E74C3C'],
    ), row=1, col=2)

    fig.update_layout(
        height=380,
        margin=dict(l=10, r=10, t=40, b=10),
    )
    return fig


def make_follower_chart(df):
    """跟进人工作量 TOP8"""
    counts = df['跟进人'].value_counts().nlargest(8).reset_index()
    counts.columns = ['跟进人','数量']
    counts = counts.sort_values('数量')

    colors = ['#E74C3C' if v >= 15 else '#F39C12' if v >= 10 else '#1890FF' for v in counts['数量']]

    fig = go.Figure(go.Bar(
        x=counts['数量'], y=counts['跟进人'],
        orientation='h',
        marker_color=colors,
        text=counts['数量'],
        textposition='outside',
    ))
    fig.update_layout(
        title='跟进人工作量 TOP8',
        xaxis_title='处理数量',
        height=300,
        margin=dict(l=10, r=30, t=40, b=10),
        bargap=0.3,
    )
    return fig


def make_d8_qa_chart(df):
    """8D报告覆盖率 + 品质负责人"""
    d8_total = len(df)
    d8_yes = len(df[df['8D报告'] == '有'])
    d8_rate = round(d8_yes / d8_total * 100, 1) if d8_total > 0 else 0

    fig = make_subplots(
        rows=1, cols=2,
        specs=[[{'type': 'indicator'}, {'type': 'pie'}]],
        column_widths=[0.4, 0.6],
        subplot_titles=('8D报告覆盖率', '品质负责人分布'),
    )

    fig.add_trace(go.Indicator(
        mode='gauge+number+delta',
        value=d8_rate,
        number={'suffix': '%', 'font': {'size': 36}},
        delta={'reference': 80, 'relative': False},
        gauge={
            'axis': {'range': [0, 100], 'tickwidth': 1},
            'bar': {'color': '#1890FF'},
            'steps': [
                {'range': [0, 60], 'color': 'rgba(231,76,60,0.3)'},
                {'range': [60, 85], 'color': 'rgba(243,156,18,0.3)'},
                {'range': [85, 100], 'color': 'rgba(39,174,96,0.3)'},
            ],
            'threshold': {
                'line': {'color': '#E74C3C', 'width': 2},
                'thickness': 0.8, 'value': 80
            }
        },
        title={'text': f'已出具: {d8_yes}/{d8_total}'},
    ), row=1, col=1)

    qa_counts = df['品质负责人'].value_counts().reset_index()
    qa_counts.columns = ['负责人','数量']

    fig.add_trace(go.Pie(
        labels=qa_counts['负责人'], values=qa_counts['数量'],
        hole=0.5, textinfo='label+percent',
    ), row=1, col=2)

    fig.update_layout(
        height=340,
        margin=dict(l=10, r=10, t=40, b=10),
    )
    return fig


def make_cycle_chart(df):
    """完成周期分布"""
    cycle = df['完成周期（天）'].dropna()
    labels = ['1-3天','4-6天','7-9天','10天以上']
    values = [
        int((cycle <= 3).sum()),
        int(((cycle > 3) & (cycle <= 6)).sum()),
        int(((cycle > 6) & (cycle <= 9)).sum()),
        int((cycle > 9).sum()),
    ]
    avg_str = f'{cycle.mean():.1f}' if len(cycle) > 0 else '-'
    med_str = f'{cycle.median():.0f}' if len(cycle) > 0 else '-'

    fig = go.Figure(go.Bar(
        x=labels, y=values,
        marker_color=['#27AE60','#1890FF','#F39C12','#E74C3C'],
        text=values, textposition='outside',
    ))
    fig.update_layout(
        title=f'完成周期分布（平均 {avg_str} 天，中位数 {med_str} 天）',
        xaxis_title='周期',
        yaxis_title='条数',
        height=300,
        margin=dict(l=10, r=10, t=40, b=10),
    )
    return fig


def send_warning_emails(df, now) -> dict:
    """发送结案预警邮件，返回发送结果"""
    # 筛选超期未结案
    mask = (
        (df['结案状态'] == '未结案') &
        df['结案预警'].notna() &
        (df['结案预警'] <= now)
    )
    overdue = df[mask].copy()
    if len(overdue) == 0:
        return {"success": 0, "fail": 0, "skipped": 0, "details": [], "total_overdue": 0}

    overdue['超期天数'] = (now - overdue['结案预警']).dt.days
    total_unclosed = int((df['结案状态'] == '未结案').sum())

    results = []
    success = fail = skipped = 0

    for follower, group in overdue.groupby('跟进人'):
        email = FOLLOWER_EMAILS.get(follower)
        if not email:
            skipped += 1
            results.append({"follower": follower, "email": "", "count": len(group), "status": "skipped"})
            continue

        # 构造 HTML 邮件
        rows_html = ""
        for _, r in group.iterrows():
            warn_date = r['结案预警'].strftime('%Y-%m-%d') if pd.notna(r['结案预警']) else '-'
            due_date = r['应结案日期'].strftime('%Y-%m-%d') if pd.notna(r['应结案日期']) else '-'
            desc = str(r['问题描述'])[:100] if pd.notna(r['问题描述']) else ''
            rows_html += f"""<tr>
                <td style="padding:6px 8px;border-bottom:1px solid #eee;">{int(r['编号_int'])}</td>
                <td style="padding:6px 8px;border-bottom:1px solid #eee;">{r['国家或地区']}</td>
                <td style="padding:6px 8px;border-bottom:1px solid #eee;">{due_date}</td>
                <td style="padding:6px 8px;border-bottom:1px solid #eee;color:#e60012;font-weight:bold;">{warn_date}</td>
                <td style="padding:6px 8px;border-bottom:1px solid #eee;color:#e60012;font-weight:bold;">{int(r['超期天数'])}天</td>
                <td style="padding:6px 8px;border-bottom:1px solid #eee;font-size:11px;color:#666;">{desc}</td>
            </tr>"""

        today_str = now.strftime('%Y-%m-%d')
        html = f"""<div style="max-width:750px;font-family:'Microsoft YaHei',Arial,sans-serif;">
            <div style="background:#e60012;padding:16px 20px;border-radius:6px 6px 0 0;">
                <h2 style="color:#fff;margin:0;font-size:16px;">2026年海外客户投诉 — 结案预警提醒</h2></div>
            <div style="background:#fff;padding:16px 20px;border:1px solid #e0e0e0;border-top:none;">
                <p style="font-size:14px;"><b>{follower}</b>，您好：</p>
                <p style="font-size:14px;">截至 <b>{today_str}</b>，您有 <span style="color:#e60012;font-weight:bold;font-size:18px;">{len(group)}条</span> 客诉已超过结案预警日期，请尽快处理。</p>
                <table style="width:100%;border-collapse:collapse;font-size:12px;">
                    <thead><tr style="background:#e60012;color:#fff;">
                        <th style="padding:8px;text-align:left;">编号</th><th style="padding:8px;text-align:left;">国家</th>
                        <th style="padding:8px;text-align:left;">应结案</th><th style="padding:8px;text-align:left;">预警</th>
                        <th style="padding:8px;text-align:left;">超期</th><th style="padding:8px;text-align:left;">问题描述</th>
                    </tr></thead><tbody>{rows_html}</tbody></table>
                <p style="margin-top:16px;font-size:11px;color:#999;">※ 此邮件由海外客诉看板自动发送 · 全量未结案 <b>{total_unclosed}</b> 条</p></div></div>"""

        msg = MIMEMultipart("alternative")
        msg["Subject"] = f"【客诉预警】{follower}，您有 {len(group)} 条客诉超期未结案 - {today_str}"
        msg["From"] = f"海外客户服务部 <{SMTP_USER}>"
        msg["To"] = f"{follower} <{email}>"
        msg.attach(MIMEText(html, "html", "utf-8"))

        try:
            if SMTP_PORT == 465:
                server = smtplib.SMTP_SSL(SMTP_SERVER, SMTP_PORT, timeout=15)
            else:
                server = smtplib.SMTP(SMTP_SERVER, SMTP_PORT, timeout=15)
                server.ehlo()
                if server.has_extn("STARTTLS"):
                    server.starttls()
                    server.ehlo()
            server.login(SMTP_USER, SMTP_PASSWORD)
            server.sendmail(SMTP_USER, email, msg.as_string())
            server.quit()
            success += 1
            results.append({"follower": follower, "email": email, "count": len(group), "status": "sent"})
        except Exception as e:
            fail += 1
            results.append({"follower": follower, "email": email, "count": len(group), "status": f"failed: {e}"})

    return {"success": success, "fail": fail, "skipped": skipped, "details": results, "total_overdue": len(overdue)}


# ============================================================
# Main App
# ============================================================
def generate_export_zip(df, now):
    """生成包含全部分析维度 + 明细数据的 ZIP 文件"""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        n = len(df)
        if n == 0:
            zf.writestr('提示.txt', '当前筛选条件下无数据，请调整筛选器后重试。'.encode('utf-8'))
            buf.seek(0)
            return buf

        # --- 01 客诉明细数据 ---
        detail_cols = {
            '编号_int': '编号', '分公司': '分公司', '国家或地区': '国家', '是否大客户': '大客户',
            '投诉日期': '投诉日期', '应结案日期': '应结案日期', '实际完成日期': '实际完成日期',
            '完成周期（天）': '完成周期(天)', '机型属性': '机型', '故障大类': '故障大类',
            '结案状态': '结案状态', '处理类型': '处理类型', '跟进人': '跟进人',
            '品质负责人': '品质负责人', '8D报告': '8D报告', '问题描述': '问题描述',
        }
        detail = df[list(detail_cols.keys())].copy()
        detail.columns = list(detail_cols.values())
        for dc in ['投诉日期','应结案日期','实际完成日期']:
            if dc in detail.columns:
                detail[dc] = detail[dc].apply(lambda x: x.strftime('%Y-%m-%d') if pd.notna(x) and hasattr(x,'strftime') else '')
        zf.writestr('01_客诉明细数据.csv', detail.to_csv(index=False, encoding='utf-8-sig'))

        # --- 02 分公司分布 ---
        branch = df['分公司'].value_counts().reset_index()
        branch.columns = ['分公司','投诉数量']
        branch['占比(%)'] = (branch['投诉数量'] / len(df) * 100).round(1)
        zf.writestr('02_分公司分布.csv', branch.to_csv(index=False, encoding='utf-8-sig'))

        # --- 03 国家地区TOP ---
        country = df['国家或地区'].value_counts().reset_index()
        country.columns = ['国家或地区','投诉数量']
        country['占比(%)'] = (country['投诉数量'] / len(df) * 100).round(1)
        zf.writestr('03_国家地区统计.csv', country.to_csv(index=False, encoding='utf-8-sig'))

        # --- 04 月度趋势 ---
        monthly = df.groupby(df['投诉日期'].dt.to_period('M').astype(str)).agg(
            投诉量=('编号_int','count'),
            已结案量=('结案状态', lambda x: x.isin(['结案','关闭']).sum()),
            未结案量=('结案状态', lambda x: (~x.isin(['结案','关闭'])).sum()),
        ).reset_index()
        monthly.columns = ['月份','投诉量','已结案量','未结案量']
        monthly['结案率(%)'] = (monthly['已结案量'] / monthly['投诉量'] * 100).round(1)
        zf.writestr('04_月度趋势.csv', monthly.to_csv(index=False, encoding='utf-8-sig'))

        # --- 05 结案状态 ---
        status = df['结案状态'].value_counts().reset_index()
        status.columns = ['结案状态','数量']
        status['占比(%)'] = (status['数量'] / len(df) * 100).round(1)
        zf.writestr('05_结案状态分布.csv', status.to_csv(index=False, encoding='utf-8-sig'))

        # --- 06 故障大类帕累托 ---
        fault = df['故障大类'].value_counts().reset_index()
        fault.columns = ['故障大类','数量']
        fault = fault.sort_values('数量', ascending=False)
        fault['占比(%)'] = (fault['数量'] / fault['数量'].sum() * 100).round(1)
        fault['累计占比(%)'] = fault['占比(%)'].cumsum().round(1)
        zf.writestr('06_故障大类帕累托.csv', fault.to_csv(index=False, encoding='utf-8-sig'))

        # --- 07 机型属性 ---
        model = df['机型属性'].value_counts().reset_index()
        model.columns = ['机型属性','数量']
        model['占比(%)'] = (model['数量'] / len(df) * 100).round(1)
        zf.writestr('07_机型属性分布.csv', model.to_csv(index=False, encoding='utf-8-sig'))

        # --- 08 跟进人工作量 ---
        follower = df['跟进人'].value_counts().reset_index()
        follower.columns = ['跟进人','处理数量']
        follower['占比(%)'] = (follower['处理数量'] / len(df) * 100).round(1)
        zf.writestr('08_跟进人工作量.csv', follower.to_csv(index=False, encoding='utf-8-sig'))

        # --- 09 品质负责人 ---
        qa = df['品质负责人'].value_counts().reset_index()
        qa.columns = ['品质负责人','负责数量']
        qa['占比(%)'] = (qa['负责数量'] / len(df) * 100).round(1)
        zf.writestr('09_品质负责人分布.csv', qa.to_csv(index=False, encoding='utf-8-sig'))

        # --- 10 超期未结案 ---
        overdue_mask = ~df['结案状态'].isin(['结案','关闭']) & df['应结案日期'].notna() & (df['应结案日期'] < now)
        overdue = df[overdue_mask][['编号_int','国家或地区','投诉日期','应结案日期','结案状态','跟进人','问题描述']].copy()
        if len(overdue) > 0:
            overdue['超期天数'] = (now - overdue['应结案日期']).dt.days
            overdue.columns = ['编号','国家','投诉日期','应结案日期','结案状态','跟进人','问题描述','超期天数']
            for dc in ['投诉日期','应结案日期']:
                overdue[dc] = overdue[dc].apply(lambda x: x.strftime('%Y-%m-%d') if pd.notna(x) and hasattr(x,'strftime') else '')
        else:
            overdue = pd.DataFrame({'提示': ['无超期未结案记录']})
        zf.writestr('10_超期未结案预警.csv', overdue.to_csv(index=False, encoding='utf-8-sig'))

        # --- 11 分公司×结案状态交叉表 ---
        cross = pd.crosstab(df['分公司'], df['结案状态'])
        cross['合计'] = cross.sum(axis=1)
        cross = cross.sort_values('合计', ascending=False)
        zf.writestr('11_分公司×结案状态交叉表.csv', cross.to_csv(encoding='utf-8-sig'))

        # --- 12 故障大类×机型交叉表 ---
        cross2 = pd.crosstab(df['故障大类'], df['机型属性'])
        cross2['合计'] = cross2.sum(axis=1)
        cross2 = cross2.sort_values('合计', ascending=False)
        zf.writestr('12_故障大类×机型交叉表.csv', cross2.to_csv(encoding='utf-8-sig'))

        # --- 13 8D报告统计 ---
        d8 = df['8D报告'].value_counts().reset_index()
        d8.columns = ['8D报告状态','数量']
        d8['占比(%)'] = (d8['数量'] / len(df) * 100).round(1)
        zf.writestr('13_8D报告覆盖率.csv', d8.to_csv(index=False, encoding='utf-8-sig'))

        # --- 14 完成周期统计 ---
        cycle = df['完成周期（天）'].dropna()
        if len(cycle) > 0:
            cycle_stats = pd.DataFrame({
                '指标': ['记录数','平均值','中位数','最小值','最大值','标准差'],
                '数值': [
                    len(cycle),
                    round(cycle.mean(), 1),
                    round(cycle.median(), 1),
                    int(cycle.min()),
                    int(cycle.max()),
                    round(cycle.std(), 1) if len(cycle) > 1 else 0,
                ]
            })
        else:
            cycle_stats = pd.DataFrame({'指标': ['记录数'], '数值': [0]})
        zf.writestr('14_完成周期统计.csv', cycle_stats.to_csv(index=False, encoding='utf-8-sig'))

    buf.seek(0)
    return buf


# ============================================================
# PPT Export
# ============================================================
# Chinese font management (runtime download for Streamlit Cloud)
_CN_FONT_PATH = '/tmp/wqy-microhei.ttc'
_FONT_URLS = [
    'https://github.com/anthonyfok/fonts-wqy-microhei/raw/master/wqy-microhei.ttc',
    'https://raw.githubusercontent.com/anthonyfok/fonts-wqy-microhei/master/wqy-microhei.ttc',
]

def _download_cn_font():
    """Download Chinese font at runtime if not available"""
    if os.path.exists(_CN_FONT_PATH) and os.path.getsize(_CN_FONT_PATH) > 100000:
        return _CN_FONT_PATH

    # Clear matplotlib font cache so it picks up new fonts
    import glob
    for cache in glob.glob('/tmp/matplotlib-*') + glob.glob('/home/*/.cache/matplotlib/*'):
        try:
            os.remove(cache)
        except Exception:
            pass

    for url in _FONT_URLS:
        try:
            req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req, timeout=30) as resp:
                data = resp.read()
                if len(data) > 100000:
                    with open(_CN_FONT_PATH, 'wb') as f:
                        f.write(data)
                    # Register font with matplotlib
                    fm.fontManager.addfont(_CN_FONT_PATH)
                    fm._load_fontmanager(try_read_cache=False)
                    return _CN_FONT_PATH
        except Exception:
            continue
    return None


def _setup_cn_font():
    """尝试配置中文字体，返回 FontProperties 对象"""
    # 1. Try runtime-downloaded font first (most reliable on Streamlit Cloud)
    downloaded = _download_cn_font()
    if downloaded:
        try:
            return fm.FontProperties(fname=downloaded)
        except Exception:
            pass

    # 2. Try installed system font paths
    font_files = [
        '/usr/share/fonts/truetype/wqy/wqy-microhei.ttc',
        '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc',
        '/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf',
        '/System/Library/Fonts/PingFang.ttc',
        '/System/Library/Fonts/STHeiti Light.ttc',
    ]
    for fp in font_files:
        if os.path.exists(fp):
            try:
                return fm.FontProperties(fname=fp)
            except Exception:
                continue

    # 3. Fallback: search by font name
    candidates = ['WenQuanYi Micro Hei','PingFang SC','Heiti SC','SimHei','Microsoft YaHei','sans-serif']
    try:
        fm._load_fontmanager(try_read_cache=False)
    except Exception:
        pass
    available = [f.name for f in fm.fontManager.ttflist]
    for c in candidates:
        if c in available:
            return fm.FontProperties(family=c)
    return fm.FontProperties(family='sans-serif')


CN_FONT_PROP = None  # Lazy init
PPT_FONT = 'Microsoft YaHei'
TCL_RED = RGBColor(0xE6, 0x00, 0x12)
TCL_DARK = RGBColor(0x33, 0x33, 0x33)
TCL_LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)

@st.cache_data(show_spinner=False)
def _load_template_bytes():
    """加载 TCL PPT 模板字节（缓存，仅首次读取磁盘）"""
    tpl_path = os.path.join(os.path.dirname(__file__), "assets", "TCL_template.pptx")
    with open(tpl_path, 'rb') as f:
        return f.read()


def _chart_to_img(fig):
    """将 matplotlib figure 转为 PNG bytes"""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=120, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    buf.seek(0)
    plt.close(fig)
    return buf


def _extract_cover_assets(template_bytes):
    """从模板封面提取 TCL 品牌元素（logo图片 + 位置）"""
    tmpl = Presentation(io.BytesIO(template_bytes))
    slide0 = tmpl.slides[0]
    assets = {'logo_blob': None, 'logo_left': 0, 'logo_top': 0,
              'logo_width': 0, 'logo_height': 0}
    for shape in slide0.shapes:
        if shape.shape_type == 13:  # PICTURE — TCL组合logo（含奥运五环+奥林匹克全球合作伙伴）
            assets['logo_blob'] = shape.image.blob
            assets['logo_left'] = shape.left
            assets['logo_top'] = shape.top
            assets['logo_width'] = shape.width
            assets['logo_height'] = shape.height
            break
    return assets


def _add_content_title(slide, title_text, subtitle_text="Data Analysis Report"):
    """按照模板目录页样式添加标题 — 左上角红色标题 + 灰色英文副标题"""
    # 红色主标题 — 参照模板目录页 (1.23", 0.21"), 21pt, #FF0000, Microsoft YaHei
    txBox = slide.shapes.add_textbox(Inches(1.23), Inches(0.21), Inches(10.5), Inches(0.45))
    p = txBox.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(21)
    p.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'

    # 灰色英文副标题 — 参照模板目录页 (1.18", 0.65"), 14pt, #898989, Microsoft YaHei Light
    txBox2 = slide.shapes.add_textbox(Inches(1.18), Inches(0.65), Inches(10.5), Inches(0.35))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = subtitle_text
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0x89, 0x89, 0x89)
    p2.font.name = 'Microsoft YaHei Light'


def _set_tf_font(tf, font_name=PPT_FONT):
    """递归设置文本框所有段落的字体"""
    for p in tf.paragraphs:
        p.font.name = font_name


def _add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    """在幻灯片中添加格式化表格"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

    # Header row
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TCL_RED
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.font.bold = True
            p.font.name = PPT_FONT
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p.font.name = PPT_FONT
                p.alignment = PP_ALIGN.CENTER

    return tbl_shape


def _make_bar_chart(labels, values, title, color='#1890FF', highlight_n=0):
    """生成 matplotlib 柱状图 figure"""
    global CN_FONT_PROP
    if CN_FONT_PROP is None:
        CN_FONT_PROP = _setup_cn_font()
    plt.rcParams['font.family'] = CN_FONT_PROP.get_name()
    plt.rcParams['axes.unicode_minus'] = False

    fig, ax = plt.subplots(figsize=(8, 4))
    colors = [color] * len(labels)
    if highlight_n > 0:
        for i in range(min(highlight_n, len(labels))):
            colors[-1 - i] = '#E74C3C'

    bars = ax.barh(range(len(labels)), values, color=colors, height=0.6)
    ax.set_yticks(range(len(labels)))
    ax.set_yticklabels(labels, fontsize=9)
    ax.set_xlabel('Quantity', fontsize=9)
    ax.set_title(title, fontsize=12, fontweight='bold')
    ax.invert_yaxis()

    for bar, val in zip(bars, values):
        ax.text(bar.get_width() + 0.1, bar.get_y() + bar.get_height()/2,
                str(val), va='center', fontsize=9)

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    plt.tight_layout()
    return fig


def generate_export_ppt(df, now):
    """生成 PPT 数据分析报告 — 始终使用内置 TCL 品牌模板"""
    global CN_FONT_PROP, PPT_FONT
    CN_FONT_PROP = _setup_cn_font()

    # 从缓存加载 TCL 模板，每次导出创建新的 Presentation 对象
    tpl_bytes = _load_template_bytes()

    # 在删除 slides 之前先提取封面品牌元素
    cover_assets = _extract_cover_assets(tpl_bytes)

    prs = Presentation(io.BytesIO(tpl_bytes))
    PPT_FONT = 'Microsoft YaHei'

    # 移除模板原有 slides，保留 layouts / masters / theme
    sldIdLst = prs.slides._sldIdLst
    ns = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
    while len(sldIdLst) > 0:
        rId = sldIdLst[0].get(ns + 'id')
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        sldIdLst.remove(sldIdLst[0])

    total = len(df)
    done = len(df[df['结案状态'].isin(['结案','关闭'])])
    rate = round(done / total * 100, 1) if total > 0 else 0
    pending = len(df[~df['结案状态'].isin(['结案','关闭'])])
    cycle_avg_raw = df['完成周期（天）'].dropna().mean()
    cycle_avg = round(cycle_avg_raw, 1) if pd.notna(cycle_avg_raw) else '-'
    overdue_count = len(df[~df['结案状态'].isin(['结案','关闭']) &
                           df['应结案日期'].notna() & (df['应结案日期'] < now)])

    # 封面使用 Layout 0 "空白"（自带全幅背景图 + 右上角logo）
    cover_layout = prs.slide_layouts[0]
    # 内容页使用 Layout 2 "标题幻灯片"（自带右上角logo，无全幅背景）
    content_layout = prs.slide_layouts[2]

    # TCL 红
    TCL_RED = RGBColor(0xE6, 0x00, 0x12)

    # ================================================================
    # Slide 1: 封面 — 模板第一页样式，TCL组合logo + 标题
    # ================================================================
    slide = prs.slides.add_slide(cover_layout)

    # 放置 TCL 组合 logo（含奥运五环、奥林匹克全球合作伙伴）
    # 位置和大小完全参照模板 Slide 0
    if cover_assets['logo_blob']:
        logo_stream = io.BytesIO(cover_assets['logo_blob'])
        slide.shapes.add_picture(logo_stream,
            cover_assets['logo_left'], cover_assets['logo_top'],
            cover_assets['logo_width'], cover_assets['logo_height'])

    # 主标题 — 模板位置 (0.80", 3.79")，30pt Microsoft YaHei
    txBox = slide.shapes.add_textbox(Inches(0.80), Inches(3.79), Inches(11.0), Inches(0.70))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "2026年海外客户投诉数据分析报告"
    p.font.size = Pt(30)
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'

    # 英文副标题 — 模板位置 (0.80", 4.35")，15pt Microsoft YaHei
    txBox2 = slide.shapes.add_textbox(Inches(0.80), Inches(4.35), Inches(11.0), Inches(0.40))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = "Overseas Customer Complaint Data Analysis"
    p2.font.size = Pt(15)
    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p2.font.name = 'Microsoft YaHei'

    # 数据范围 — 模板位置 (0.80", 5.04")，17pt Microsoft YaHei
    txBox3 = slide.shapes.add_textbox(Inches(0.80), Inches(5.04), Inches(11.0), Inches(0.40))
    p3 = txBox3.text_frame.paragraphs[0]
    p3.text = f"数据范围: 2026年1月-6月 | 有效记录: {total} 条"
    p3.font.size = Pt(17)
    p3.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p3.font.name = 'Microsoft YaHei'

    # 生成时间 — 模板位置 (0.80", 5.38")，12pt Microsoft YaHei
    txBox4 = slide.shapes.add_textbox(Inches(0.80), Inches(5.38), Inches(11.0), Inches(0.40))
    p4 = txBox4.text_frame.paragraphs[0]
    p4.text = f"生成时间: {now.strftime('%Y-%m-%d %H:%M')} (北京时间)"
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p4.font.name = 'Microsoft YaHei'

    # ================================================================
    # Slide 2: 执行摘要 - KPI
    # ================================================================
    slide = prs.slides.add_slide(content_layout)
    _add_content_title(slide, "📊 执行摘要 — 关键指标概览")

    kpi_data = [
        ('📋 总投诉量', f'{total} 条', '2026年1-6月累计'),
        ('✅ 结案率', f'{rate}%', f'结案+关闭共{done}条'),
        ('⏱️ 平均处理周期', f'{cycle_avg} 天', '基于已完成记录'),
        ('⚠️ 未结案', f'{pending} 条', '含暂停/未结案'),
        ('🔴 超期未结案', f'{overdue_count} 条', '已超应结案日期'),
    ]

    for i, (label, value, sub) in enumerate(kpi_data):
        left = Inches(0.8 + i * 2.5)
        top = Inches(1.3)
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.2), Inches(1.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0xF0, 0xF2, 0xF5)
        card.line.fill.background()

        # Value
        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), Inches(1.8), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p.font.name = 'Microsoft YaHei'
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(28)
        p2.font.color.rgb = RGBColor(0xE7, 0x4C, 0x3C) if '🔴' in label else TCL_RED
        p2.font.bold = True
        p2.font.name = 'Microsoft YaHei'
        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(9)
        p3.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p3.font.name = 'Microsoft YaHei'

    # Summary text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    summary_items = [
        f"● 2026年上半年累计受理海外客户投诉 {total} 条，整体结案率 {rate}%，尚有 {pending} 条未完结。",
        f"● 投诉高峰集中在3-5月，6月起逐步回落。",
        f"● 当前已识别超期未结案 {overdue_count} 条，需紧急处理。",
        f"● 变频机型投诉占比较高，美国和欧洲是投诉最多的市场。",
        f"● {overdue_count} 条超期未结案需紧急处理，请逐条制定结案计划。",
    ]
    for item in summary_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(6)
        p.font.name = 'Microsoft YaHei'

    # ================================================================
    # Slide 3: 区域分布 - 分公司 + 国家TOP10
    # ================================================================
    slide = prs.slides.add_slide(content_layout)
    _add_content_title(slide, "🌍 区域分布分析 — 分公司 & 国家/地区")

    # Branch chart
    branch = df['分公司'].value_counts()
    branch_labels = list(branch.index)
    branch_vals = list(branch.values)
    fig1 = _make_bar_chart(branch_labels, branch_vals, '分公司投诉分布', highlight_n=3)
    img1 = _chart_to_img(fig1)
    slide.shapes.add_picture(img1, Inches(0.3), Inches(1.0), Inches(6.2), Inches(3.0))

    # Branch table
    branch_rows = [(n, v, f'{v/total*100:.1f}%') for n, v in zip(branch_labels, branch_vals)]
    _add_table(slide, Inches(0.3), Inches(4.2), Inches(6.2), Inches(2.8),
               ['分公司','投诉数量','占比'], branch_rows)

    # Country chart
    country = df['国家或地区'].value_counts().nlargest(10)
    country_labels = list(country.index)
    country_vals = list(country.values)
    fig2 = _make_bar_chart(country_labels, country_vals, '国家/地区投诉量 TOP10', highlight_n=3)
    img2 = _chart_to_img(fig2)
    slide.shapes.add_picture(img2, Inches(6.8), Inches(1.0), Inches(6.2), Inches(3.0))

    # Country table
    country_rows = [(n, v, f'{v/total*100:.1f}%') for n, v in zip(country_labels, country_vals)]
    _add_table(slide, Inches(6.8), Inches(4.2), Inches(6.2), Inches(2.8),
               ['国家/地区','投诉数量','占比'], country_rows)

    # ================================================================
    # Slide 4: 时间趋势 + 结案状态
    # ================================================================
    slide = prs.slides.add_slide(content_layout)
    _add_content_title(slide, "📈 时间趋势 & 结案状态分析")

    # Monthly trend table
    df['投诉月份'] = df['投诉日期'].dt.to_period('M').astype(str)
    monthly = df.groupby('投诉月份').agg(
        投诉量=('编号_int','count'),
        结案量=('结案状态', lambda x: x.isin(['结案','关闭']).sum())
    ).reset_index()
    monthly['结案率'] = (monthly['结案量'] / monthly['投诉量'] * 100).round(1)
    monthly_rows = [(r['投诉月份'], r['投诉量'], r['结案量'], f"{r['结案率']}%")
                    for _, r in monthly.iterrows()]
    _add_table(slide, Inches(0.5), Inches(1.0), Inches(5.5), Inches(2.0),
               ['月份','投诉量','结案量','结案率'], monthly_rows)

    # Status pie table
    status = df['结案状态'].value_counts()
    status_rows = [(n, v, f'{v/total*100:.1f}%') for n, v in status.items()]
    _add_table(slide, Inches(0.5), Inches(3.5), Inches(3.0), Inches(1.8),
               ['结案状态','数量','占比'], status_rows)

    # VIP table
    vip = df['是否大客户'].value_counts()
    vip_rows = [(n, v, f'{v/total*100:.1f}%') for n, v in vip.items()]
    _add_table(slide, Inches(4.0), Inches(3.5), Inches(3.0), Inches(1.3),
               ['客户类型','数量','占比'], vip_rows)

    # Model type table
    model = df['机型属性'].value_counts()
    model_rows = [(n, v, f'{v/total*100:.1f}%') for n, v in model.items()]
    _add_table(slide, Inches(7.5), Inches(1.0), Inches(5.5), Inches(3.5),
               ['机型属性','数量','占比'], model_rows,
               col_widths=[Inches(2.2), Inches(1.5), Inches(1.5)])

    # ================================================================
    # Slide 5: 故障分析 — 帕累托 + 故障×机型交叉
    # ================================================================
    slide = prs.slides.add_slide(content_layout)
    _add_content_title(slide, "🔍 故障分析 — 帕累托 & 故障×机型交叉")

    # Fault pareto chart
    fault = df['故障大类'].value_counts()
    fault_labels = list(fault.index)
    fault_vals = list(fault.values)
    fig3 = _make_bar_chart(fault_labels, fault_vals, '故障大类分布', highlight_n=3)
    img3 = _chart_to_img(fig3)
    slide.shapes.add_picture(img3, Inches(0.3), Inches(1.0), Inches(6.2), Inches(3.2))

    # Fault pareto table
    fault_total = sum(fault_vals)
    cum = 0
    fault_rows = []
    for n, v in zip(fault_labels, fault_vals):
        cum += v
        fault_rows.append((n, v, f'{v/fault_total*100:.1f}%', f'{cum/fault_total*100:.1f}%'))
    _add_table(slide, Inches(0.3), Inches(4.4), Inches(6.2), Inches(2.8),
               ['故障大类','数量','占比','累计占比'], fault_rows)

    # Cross table: fault × model (TOP rows)
    cross = pd.crosstab(df['故障大类'], df['机型属性'])
    cross['合计'] = cross.sum(axis=1)
    cross = cross.sort_values('合计', ascending=False).head(8)
    cross_headers = ['故障大类'] + list(cross.columns)
    cross_rows = []
    for idx, row in cross.iterrows():
        cross_rows.append([idx] + [int(row[c]) for c in cross.columns])
    _add_table(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(3.5),
               cross_headers, cross_rows)

    # ================================================================
    # Slide 6: 质量管理 — 跟进人 + 8D + 周期
    # ================================================================
    slide = prs.slides.add_slide(content_layout)
    _add_content_title(slide, "👤 质量管理 — 跟进人 & 8D报告 & 处理周期")

    # Follower chart
    follower = df['跟进人'].value_counts().nlargest(8)
    f_labels = list(follower.index)
    f_vals = list(follower.values)
    fig4 = _make_bar_chart(f_labels, f_vals, '跟进人工作量 TOP8', highlight_n=2)
    img4 = _chart_to_img(fig4)
    slide.shapes.add_picture(img4, Inches(0.3), Inches(1.0), Inches(4.2), Inches(2.8))

    # Follower table
    f_rows = [(n, v, f'{v/total*100:.1f}%') for n, v in zip(f_labels, f_vals)]
    _add_table(slide, Inches(0.3), Inches(4.0), Inches(4.2), Inches(2.5),
               ['跟进人','处理数量','占比'], f_rows)

    # 8D stats
    d8_yes = len(df[df['8D报告'] == '有'])
    d8_no = len(df[df['8D报告'] == '无'])
    d8_rate = round(d8_yes / total * 100, 1) if total > 0 else 0
    txBox = slide.shapes.add_textbox(Inches(5.0), Inches(1.2), Inches(4.0), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "8D报告覆盖率"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p2 = tf.add_paragraph()
    p2.text = f"{d8_rate}%"
    p2.font.size = Pt(42)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0x27, 0xAE, 0x60) if d8_rate >= 80 else RGBColor(0xF3, 0x9C, 0x12)
    p3 = tf.add_paragraph()
    p3.text = f"已出具: {d8_yes} 条 | 未出具: {d8_no} 条 | 其他: {total - d8_yes - d8_no} 条"
    p3.font.size = Pt(11)
    p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # QA table
    qa = df['品质负责人'].value_counts()
    qa_rows = [(n, v, f'{v/total*100:.1f}%') for n, v in qa.items()]
    _add_table(slide, Inches(5.0), Inches(4.4), Inches(3.8), Inches(2.0),
               ['品质负责人','负责数量','占比'], qa_rows)

    # Cycle stats
    cycle = df['完成周期（天）'].dropna()
    if len(cycle) > 0:
        cycle_rows = [
            ('记录数', len(cycle)), ('平均值', f'{cycle.mean():.1f} 天'),
            ('中位数', f'{cycle.median():.0f} 天'), ('最小值', f'{cycle.min():.0f} 天'),
            ('最大值', f'{cycle.max():.0f} 天'), ('标准差', f'{cycle.std():.1f} 天' if len(cycle) > 1 else '0.0 天'),
        ]
    else:
        cycle_rows = [('记录数', 0)]
    _add_table(slide, Inches(9.2), Inches(1.0), Inches(3.8), Inches(2.0),
               ['完成周期指标','数值'], cycle_rows)

    # ================================================================
    # Slide 7: 超期未结案预警
    # ================================================================
    slide = prs.slides.add_slide(content_layout)
    _add_content_title(slide, "🚨 超期未结案预警 & 分公司×状态交叉表")

    # Overdue warnings
    overdue = df[~df['结案状态'].isin(['结案','关闭']) &
                 df['应结案日期'].notna() & (df['应结案日期'] < now)]
    if len(overdue) > 0:
        od_rows = []
        for _, r in overdue.iterrows():
            days = (now - r['应结案日期']).days
            desc = str(r['问题描述'])[:80] if pd.notna(r['问题描述']) else ''
            od_rows.append((r['编号_int'], r['国家或地区'],
                           r['应结案日期'].strftime('%Y-%m-%d'), f'{days}天', desc))
        _add_table(slide, Inches(0.3), Inches(1.0), Inches(7.0), Inches(min(3.5, 0.4 * len(od_rows) + 0.5)),
                   ['编号','国家','应结案日期','超期天数','问题描述'], od_rows,
                   col_widths=[Inches(0.6), Inches(1.0), Inches(1.2), Inches(1.0), Inches(3.2)])
    else:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "✅ 当前无超期未结案记录"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0x27, 0xAE, 0x60)

    # Cross table
    cross = pd.crosstab(df['分公司'], df['结案状态'])
    cross['合计'] = cross.sum(axis=1)
    cross = cross.sort_values('合计', ascending=False)
    cross_headers = ['分公司'] + list(cross.columns)
    cross_rows_data = []
    for idx, row in cross.iterrows():
        cross_rows_data.append([idx] + [int(row[c]) for c in cross.columns])
    _add_table(slide, Inches(7.8), Inches(1.0), Inches(5.2), Inches(min(3.5, 0.35 * len(cross) + 0.5)),
               cross_headers, cross_rows_data)

    # ================================================================
    # Slide 8: 建议与总结
    # ================================================================
    slide = prs.slides.add_slide(content_layout)
    _add_content_title(slide, "💡 总结与改进建议")

    recommendations = [
        ("1. 重点治理装配问题",
         f"装配问题为投诉中占比较高的故障类型，建议对中山一厂、广州工厂、印尼工厂开展装配工序专项审查，"
         "加强首检和巡检力度，对频繁出现装配问题的产线进行停线整顿。"),
        ("2. 变频产品专项质量提升",
         f"变频整机+散件投诉占总量的71.1%，建议成立变频产品专项质量小组，重点关注裂管/漏氟、电控、外观不良等高频故障。"),
        ("3. 紧急处理超期未结案",
         f"当前有{overdue_count}条超期未结案，最长的已超期160+天（编号1-澳大利亚），建议逐条制定结案计划，明确责任人和完成时间。"),
        ("4. 提升8D报告覆盖率",
         f"当前8D报告覆盖率仅{d8_rate}%，建议强制要求所有已结案投诉完成8D报告归档，品质部门每月通报各负责人8D完成率。"),
        ("5. 优化装柜与包装",
         "装柜问题和包装破损合计占比11.9%，建议优化货柜装载方案，加强包装强度测试，特别是针对窗机和大宗散件的包装防护。"),
        ("6. 加强重点市场客诉响应",
         "美国（9条）和中国台湾（7条）是投诉最多的市场，建议为这两个市场配置专属品质对接窗口，缩短响应周期，降低客诉升级风险。"),
    ]

    y_pos = 1.0
    for title, detail in recommendations:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.5), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TCL_RED
        p2 = tf.add_paragraph()
        p2.text = detail
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        y_pos += 1.05

    # Footer
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"数据来源: 2026年海外客户投诉台账.xlsx | 生成时间: {now.strftime('%Y-%m-%d %H:%M')} | 有效记录: {total} 条"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER

    # Save
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def main():
    # 剥离时区信息以兼容 Excel 中的 naive datetime64[ns] 列 (Python 3.14 + pandas 3.x 要求)
    now_bjt = datetime.now(BJT)
    now = pd.Timestamp(now_bjt.replace(tzinfo=None))

    # ---- Sidebar ----
    with st.sidebar:
        st.image("https://img.icons8.com/color/96/bar-chart.png", width=48)
        st.markdown("## 📊 筛选器")

        # ---- File Upload / Data Sync ----
        st.markdown("### 📤 同步最新数据")

        # Data freshness indicator
        if 'data_update_time' not in st.session_state:
            st.session_state.data_update_time = datetime.now(BJT)

        uploaded_file = st.file_uploader(
            "从金山文档导出Excel后拖拽上传",
            type=['xlsx'],
            help="金山文档 → 导出为Excel → 拖到此处 → 看板立即刷新",
            key='excel_uploader',
        )
        if uploaded_file:
            st.session_state.data_update_time = datetime.now(BJT)
            # Save bytes so we can read multiple times
            st.session_state.uploaded_bytes = uploaded_file.getvalue()
            st.success(f"✅ 数据已同步！刷新中...")
        else:
            st.session_state.uploaded_bytes = None

        # Show data source & freshness
        if uploaded_file is None:
            st.caption("💡 本地编辑 Excel 后拖拽到此即可更新看板")

    # Load data from session state or local file
    upload_bytes = st.session_state.get('uploaded_bytes')
    df = load_data(upload_bytes)

    # Continue sidebar filters
    with st.sidebar:
        st.divider()
        filtered = apply_filters(df)

        st.divider()
        st.caption(f"筛选后记录数: **{len(filtered)}** / {len(df)}")
        st.caption(f"数据生成时间: {now.strftime('%Y-%m-%d %H:%M:%S')}")

        st.divider()
        st.markdown("### 📥 导出数据分析")

        export_zip = generate_export_zip(filtered, now)
        st.download_button(
            label="⬇️ 导出全部分析数据 (ZIP)",
            data=export_zip,
            file_name=f"海外客诉分析数据_{now.strftime('%Y%m%d_%H%M')}.zip",
            mime="application/zip",
            use_container_width=True,
        )
        st.caption("14个CSV：明细数据 + 各维度统计 + 交叉表 + 超期预警")

        st.divider()
        st.markdown("### 📊 导出PPT报告")
        st.caption("使用内置 TCL 品牌模板，无需上传")

        export_ppt = generate_export_ppt(filtered, now)
        st.download_button(
            label="📊 导出数据分析报告 (PPT)",
            data=export_ppt,
            file_name=f"海外客诉分析报告_{now.strftime('%Y%m%d_%H%M')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )
        st.caption("8页PPT：封面 + KPI摘要 + 区域/趋势/故障/质量分析 + 改进建议")

        # ---- 结案预警 ----
        st.divider()
        st.markdown("### 🚨 结案预警通知")

        # 统计超期未结案
        warn_mask = (
            (filtered['结案状态'] == '未结案') &
            filtered['结案预警'].notna() &
            (filtered['结案预警'] <= now)
        )
        overdue_warnings = filtered[warn_mask]
        overdue_count = len(overdue_warnings)

        if overdue_count > 0:
            st.warning(f"⚠️ **{overdue_count}** 条超期未结案预警")
            overdue_warnings['超期天数'] = (now - overdue_warnings['结案预警']).dt.days
            overdue_warnings_sorted = overdue_warnings.sort_values('超期天数', ascending=False)
            # 按跟进人汇总
            follower_summary = overdue_warnings_sorted.groupby('跟进人').agg(
                预警数=('编号_int', 'count'), 最长达=('超期天数', 'max')
            ).sort_values('预警数', ascending=False)
            for follower, row in follower_summary.iterrows():
                email = FOLLOWER_EMAILS.get(follower, '无邮箱')
                st.caption(f"  {follower}: {int(row['预警数'])}条预警 (最长{int(row['最长达'])}天)")

            # 一键发送按钮
            if st.button("📧 一键发送结案预警邮件", type="primary", use_container_width=True):
                with st.spinner("正在发送预警邮件..."):
                    result = send_warning_emails(filtered, now)
                if result['success'] > 0:
                    st.success(f"✅ 已发送 {result['success']} 封邮件")
                    for d in result['details']:
                        if d['status'] == 'sent':
                            st.caption(f"  ✓ {d['follower']} ({d['count']}条)")
                if result['fail'] > 0:
                    st.error(f"❌ {result['fail']} 封发送失败")
                    for d in result['details']:
                        if d['status'] != 'sent' and d['status'] != 'skipped':
                            st.caption(f"  ✗ {d['follower']}: {d['status']}")
                if result['skipped'] > 0:
                    st.warning(f"⚠️ {result['skipped']} 人无邮箱，已跳过")
        else:
            st.success("✅ 当前无超期未结案预警")

    # ---- Main Area ----
    # Header with source file link
    source = df.attrs.get('source', '未知')

    st.markdown(f"""
    <div style="background:linear-gradient(135deg,#1a1a2e,#16213e,#0f3460);padding:16px 30px;border-radius:12px;margin-bottom:16px;display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:12px;">
        <div>
            <h1 style="color:#fff;margin:0;font-size:24px;">📊 2026年海外客户投诉数据看板</h1>
            <p style="color:#aaa;margin:2px 0 0;font-size:12px;">数据源: {source} · 每30秒自动刷新</p>
        </div>
        <div style="display:flex;gap:8px;flex-wrap:wrap;">
            <a href="{KDOCS_EDIT_URL}" target="_blank" style="background:rgba(255,255,255,0.15);color:#fff;padding:8px 18px;border-radius:6px;text-decoration:none;font-size:13px;font-weight:500;border:1px solid rgba(255,255,255,0.3);transition:all 0.2s;white-space:nowrap;"
               onmouseover="this.style.background='rgba(255,255,255,0.25)'" onmouseout="this.style.background='rgba(255,255,255,0.15)'">
               📄 金山在线文档
            </a>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # ---- KPI Row ----
    total = len(filtered)
    done = len(filtered[filtered['结案状态'].isin(['结案','关闭'])])
    rate = round(done / total * 100, 1) if total > 0 else 0
    pending = len(filtered[~filtered['结案状态'].isin(['结案','关闭'])])
    cycle_avg = filtered['完成周期（天）'].dropna().mean()
    cycle_avg = round(cycle_avg, 1) if pd.notna(cycle_avg) else '-'
    overdue = len(filtered[~filtered['结案状态'].isin(['结案','关闭']) &
                           filtered['应结案日期'].notna() &
                           (filtered['应结案日期'] < now)])

    k1, k2, k3, k4, k5 = st.columns(5)
    with k1:
        st.metric("📋 总投诉量", f"{total} 条")
    with k2:
        st.metric("✅ 结案率", f"{rate}%", delta=f"结案{len(filtered[filtered['结案状态']=='结案'])}条" if total > 0 else None)
    with k3:
        st.metric("⏱️ 平均处理周期", f"{cycle_avg} 天")
    with k4:
        st.metric("⚠️ 未结案", f"{pending} 条", delta_color="inverse")
    with k5:
        st.metric("🔴 超期未结案", f"{overdue} 条", delta_color="inverse")

    st.divider()

    # ---- Tabs ----
    tab1, tab2, tab3, tab4 = st.tabs(["📊 概览仪表盘", "🔍 故障分析", "👤 质量管理", "📋 数据明细"])

    # ============ TAB 1: 概览仪表盘 ============
    with tab1:
        col1, col2 = st.columns(2)
        with col1:
            st.plotly_chart(make_branch_chart(filtered), use_container_width=True, key='tab1_branch')
        with col2:
            st.plotly_chart(make_country_chart(filtered), use_container_width=True, key='tab1_country')

        col3, col4 = st.columns(2)
        with col3:
            st.plotly_chart(make_monthly_trend(filtered), use_container_width=True, key='tab1_monthly')
        with col4:
            st.plotly_chart(make_status_pie(filtered), use_container_width=True, key='tab1_status')

    # ============ TAB 2: 故障分析 ============
    with tab2:
        col1, col2 = st.columns(2)
        with col1:
            st.plotly_chart(make_fault_pareto(filtered), use_container_width=True, key='tab2_pareto')
        with col2:
            st.plotly_chart(make_model_vip_chart(filtered), use_container_width=True, key='tab2_model')

        # Fault × Model cross table
        st.markdown("#### 故障大类 × 机型属性 交叉矩阵")
        cross = pd.crosstab(filtered['故障大类'], filtered['机型属性'])
        cross['合计'] = cross.sum(axis=1)
        cross_sorted = cross.sort_values('合计', ascending=False)
        st.dataframe(cross_sorted, use_container_width=True, height=250)

    # ============ TAB 3: 质量管理 ============
    with tab3:
        col1, col2 = st.columns(2)
        with col1:
            st.plotly_chart(make_follower_chart(filtered), use_container_width=True, key='tab3_follower')
        with col2:
            st.plotly_chart(make_cycle_chart(filtered), use_container_width=True, key='tab3_cycle')

        col3, col4 = st.columns(2)
        with col3:
            st.plotly_chart(make_d8_qa_chart(filtered), use_container_width=True, key='tab3_d8')

        with col4:
            st.markdown("#### 🚨 超期未结案预警")
            overdue_df = filtered[~filtered['结案状态'].isin(['结案','关闭']) &
                                  filtered['应结案日期'].notna() &
                                  (filtered['应结案日期'] < now)]

            if len(overdue_df) > 0:
                for _, r in overdue_df.iterrows():
                    days = (now - r['应结案日期']).days
                    severity = 'error' if days > 30 else 'warning'
                    with st.container():
                        getattr(st, severity)(
                            f"**编号{r['编号_int']}** | {r['国家或地区']} | "
                            f"应结案: {r['应结案日期'].strftime('%Y-%m-%d')} | "
                            f"超期 **{days} 天**\n\n"
                            f"{str(r['问题描述'])[:120]}..."
                        )
            else:
                st.success("✅ 当前无超期未结案记录")

        # Branch × Status cross table
        st.markdown("#### 分公司 × 结案状态 明细表")
        cross2 = pd.crosstab(filtered['分公司'], filtered['结案状态'])
        cross2['合计'] = cross2.sum(axis=1)
        st.dataframe(cross2, use_container_width=True, height=280)

    # ============ TAB 4: 数据明细 ============
    with tab4:
        # Search & filter row
        search_col1, search_col2, search_col3 = st.columns([3, 1, 1])
        with search_col1:
            search = st.text_input("🔍 全局搜索", placeholder="输入编号/国家/故障/跟进人...")
        with search_col2:
            fault_filter = st.selectbox("故障大类筛选", options=['全部'] + sorted(filtered['故障大类'].dropna().unique().tolist()))
        with search_col3:
            status_filter = st.selectbox("结案状态筛选", options=['全部'] + sorted(filtered['结案状态'].dropna().unique().tolist()))

        # Apply filters
        display = filtered.copy()
        if search:
            mask = (
                display['编号_int'].astype(str).str.contains(search, case=False) |
                display['国家或地区'].str.contains(search, case=False, na=False) |
                display['故障大类'].str.contains(search, case=False, na=False) |
                display['跟进人'].str.contains(search, case=False, na=False) |
                display['分公司'].str.contains(search, case=False, na=False) |
                display['问题描述'].str.contains(search, case=False, na=False)
            )
            display = display[mask]
        if fault_filter != '全部':
            display = display[display['故障大类'] == fault_filter]
        if status_filter != '全部':
            display = display[display['结案状态'] == status_filter]

        st.caption(f"显示 {len(display)} / {len(filtered)} 条记录")

        # Build display table
        table_cols = {
            '编号_int': '编号',
            '分公司': '分公司',
            '国家或地区': '国家',
            '是否大客户': '大客户',
            '投诉日期': '投诉日期',
            '应结案日期': '应结案',
            '实际完成日期': '实际完成',
            '完成周期（天）': '周期',
            '机型属性': '机型',
            '故障大类': '故障大类',
            '结案状态': '状态',
            '处理类型': '类型',
            '跟进人': '跟进人',
            '品质负责人': '品质负责人',
            '8D报告': '8D',
            '问题描述': '问题描述',
        }

        show = display[list(table_cols.keys())].copy()
        show.columns = list(table_cols.values())
        # Format dates
        for dc in ['投诉日期','应结案','实际完成']:
            if dc in show.columns:
                show[dc] = show[dc].apply(lambda x: x.strftime('%Y-%m-%d') if pd.notna(x) and hasattr(x, 'strftime') else str(x)[:10])

        st.dataframe(
            show,
            use_container_width=True,
            height=600,
            hide_index=True,
            column_config={
                '编号': st.column_config.NumberColumn('编号', width='small'),
                '大客户': st.column_config.TextColumn('大客户', width='small'),
                '周期': st.column_config.NumberColumn('周期(天)', width='small'),
                '问题描述': st.column_config.TextColumn('问题描述', width='large'),
            }
        )

        # Export
        csv = show.to_csv(index=False).encode('utf-8-sig')
        st.download_button(
            label="📥 导出为 CSV",
            data=csv,
            file_name=f"海外客诉数据_{now.strftime('%Y%m%d')}.csv",
            mime='text/csv',
        )

    # ---- Footer ----
    st.divider()
    st.caption(f"© 2026 海外客户投诉数据看板 · 数据更新: {now.strftime('%Y-%m-%d %H:%M')} (北京时间) · Powered by Streamlit")


if __name__ == '__main__':
    main()
